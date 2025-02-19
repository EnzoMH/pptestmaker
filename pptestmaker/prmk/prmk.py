from .layout import add_header_title, add_page_number
from .section import calculate_default_proposal_pages, calculate_proposal_summary_pages, calculate_detailed_content, recommend_pages
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import math

def create_presentation(slide_count, save_path, head_title, subtitle, section0, section1, subsections1, section2, subsections2, section3, subsections3, section4, subsections4, add_additional_page):
    prs = Presentation()
    total_slides = math.ceil(slide_count * 1.1)

    default_pages, cover_page, table_of_contents_page, last_page, project_plan_pages, additional_proposal_pages = calculate_default_proposal_pages()
    summary_pages = calculate_proposal_summary_pages(total_slides)
    preview_pages, full_content_pages = calculate_detailed_content(total_slides, default_pages, summary_pages)

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title.text = head_title
    subtitle_shape.text = subtitle

    set_font(title.text_frame, "Malgun Gothic", Pt(44))
    set_font(subtitle_shape.text_frame, "Malgun Gothic", Pt(24))

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Cm(0), Cm(0), Cm(25.4), Cm(3))
    text_frame = title.text_frame
    text_frame.text = section0
    set_font(title, "Malgun Gothic", Pt(36), bold=True)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    section1_start = 2
    section1_end = section1_start + summary_pages - 1
    section2_start = section1_end + 1
    section2_end = section2_start + preview_pages - 1
    section3_start = section2_end + 1
    section3_end = section3_start + full_content_pages - 1
    section4_start = section3_end + 1
    section4_end = section4_start + project_plan_pages - 1

    add_section(slide, f"Ⅰ. {section1}", section1_start, section1_end, 3)
    for i, subsection in enumerate(subsections1, start=1):
        add_section(slide, f"Ⅰ-{i}. {subsection}", section1_start + i, section1_start + i, 3 + i)

    add_section(slide, f"Ⅱ. {section2}", section2_start, section2_end, 4 + len(subsections1))
    for i, subsection in enumerate(subsections2, start=1):
        add_section(slide, f"Ⅱ-{i}. {subsection}", section2_start + i, section2_start + i, 4 + len(subsections1) + i)

    add_section(slide, f"Ⅲ. {section3}", section3_start, section3_end, 5 + len(subsections1) + len(subsections2))
    for i, subsection in enumerate(subsections3, start=1):
        add_section(slide, f"Ⅲ-{i}. {subsection}", section3_start + i, section3_start + i, 5 + len(subsections1) + len(subsections2) + i)

    add_section(slide, f"Ⅳ. {section4}", section4_start, section4_end, 6 + len(subsections1) + len(subsections2) + len(subsections3))
    for i, subsection in enumerate(subsections4, start=1):
        add_section(slide, f"Ⅳ-{i}. {subsection}", section4_start + i, section4_start + i, 6 + len(subsections1) + len(subsections2) + len(subsections3) + i)

    if add_additional_page == 0:
        add_section(slide, "Ⅴ. 추가제안", section4_end + 1, section4_end + 1, 7 + len(subsections1) + len(subsections2) + len(subsections3) + len(subsections4))

    add_page_number(slide, 1)

    slide_layout = prs.slide_layouts[6]
    page_number = 2
    slide = prs.slides.add_slide(slide_layout)
    add_header_title(slide, f"Ⅰ. {section1}")
    add_page_number(slide, page_number)
    page_number += 1

    for _ in range(summary_pages - 1):
        slide = prs.slides.add_slide(slide_layout)
        add_header_title(slide, f"Ⅰ. {section1}")
        add_page_number(slide, page_number)
        page_number += 1

    for subsection in subsections1:
        slide = prs.slides.add_slide(slide_layout)
        add_header_title(slide, subsection)
        add_page_number(slide, page_number)
        page_number += 1

    slide = prs.slides.add_slide(slide_layout)
    add_header_title(slide, f"Ⅱ. {section2}")
    add_page_number(slide, page_number)
    page_number += 1

    for subsection in subsections2:
        slide = prs.slides.add_slide(slide_layout)
        add_header_title(slide, subsection)
        add_page_number(slide, page_number)
        page_number += 1

    slide = prs.slides.add_slide(slide_layout)
    add_header_title(slide, f"Ⅲ. {section3}")
    add_page_number(slide, page_number)
    page_number += 1

    for subsection in subsections3:
        slide = prs.slides.add_slide(slide_layout)
        add_header_title(slide, subsection)
        add_page_number(slide, page_number)
        page_number += 1

    slide = prs.slides.add_slide(slide_layout)
    add_header_title(slide, f"Ⅳ. {section4}")
    add_page_number(slide, page_number)
    page_number += 1

    for subsection in subsections4:
        slide = prs.slides.add_slide(slide_layout)
        add_header_title(slide, subsection)
        add_page_number(slide, page_number)
        page_number += 1

    if add_additional_page == 0:
        slide = prs.slides.add_slide(slide_layout)
        add_header_title(slide, "Ⅴ. 추가제안")
        add_page_number(slide, page_number)
        page_number += 1

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Cm(0), Cm(8.46), Cm(25.4), Cm(2.54))
    text_frame = title.text_frame
    text_frame.text = "감사합니다"
    set_font(title, "Malgun Gothic", Pt(44))
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_shape = slide.shapes.add_textbox(Cm(0), Cm(18.2), Cm(25.4), Cm(0.86))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "주식회사 지아이웍스"
    set_font(subtitle_frame, "Malgun Gothic", Pt(14), color=RGBColor(128, 128, 128))
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    output_file = os.path.join(save_path, 'proposal.pptx')
    prs.save(output_file)

    if os.path.exists(output_file):
        print(f"파일이 성공적으로 생성되었습니다: {output_file}")
    else:
        print("파일 생성에 실패했습니다.")

    print(f"파일 저장 경로: {output_file}")
