from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN  # PP_ALIGN 추가
from pptx.dml.color import RGBColor
from .utils import set_font

def add_header_title(slide, text, font_size=7):
    header = slide.shapes.add_textbox(Cm(23.5), Cm(0.2), Cm(1.9), Cm(0.6))
    text_frame = header.text_frame
    text_frame.text = text
    set_font(header, "Malgun Gothic", Pt(font_size), bold=True)
    text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_page_number(slide, page_number):
    page_number_box = slide.shapes.add_textbox(Cm(24.4), Cm(18.53), Cm(1.0), Cm(0.5))
    text_frame = page_number_box.text_frame
    text_frame.text = f"- {page_number} -"
    set_font(page_number_box, "Malgun Gothic", Pt(6), bold=True)
    text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_section(slide, text, page_start, page_end, top):
    content = slide.shapes.add_textbox(Cm(1), Cm(top), Cm(23.4), Cm(1))
    text_frame = content.text_frame
    text_frame.text = f"{text} (p. {page_start}-{page_end})"
    set_font(content, "Malgun Gothic", Pt(10))
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_subsection(slide, text, top):
    content = slide.shapes.add_textbox(Cm(2), Cm(top), Cm(22.4), Cm(0.6))
    text_frame = content.text_frame
    text_frame.text = text
    set_font(content, "Malgun Gothic", Pt(8))
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
